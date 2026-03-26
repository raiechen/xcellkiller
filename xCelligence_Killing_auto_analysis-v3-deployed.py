"""
app.py  –  PPT, Excel & Word → Markdown Converter (Streamlit)
--------------------------------------------------------------
Upload a .pptx, .xlsx/.xls, or .docx file and download a clean
Markdown file with all text, headings, bullet points, and tables.

Deploy to Streamlit Cloud:
    1. Push this repo to GitHub
    2. Go to share.streamlit.io → New app → select repo
    3. Set "Main file path" to  app.py
"""

import io
import os
import streamlit as st
import pandas as pd
import pdfplumber
from pptx import Presentation
from docx import Document
from docx.oxml.ns import qn
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph

# ─────────────────────────────────────────────
# Page config (must be first Streamlit call)
# ─────────────────────────────────────────────
st.set_page_config(
    page_title="Docs → Markdown Converter",
    page_icon="📑",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ─────────────────────────────────────────────
# Custom CSS – dark gradient theme
# ─────────────────────────────────────────────
st.markdown(
    """
    <style>
    /* ── Global ──────────────────────────────── */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap');

    html, body, [class*="css"] {
        font-family: 'Inter', sans-serif;
    }

    .stApp {
        background: #f5f7fa;
        min-height: 100vh;
    }

    /* ── Hero banner ─────────────────────────── */
    .hero {
        text-align: center;
        padding: 2.5rem 1rem 1.5rem;
    }
    .hero h1 {
        font-size: 2.8rem;
        font-weight: 700;
        background: linear-gradient(90deg, #7c3aed, #2563eb, #0891b2);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 0.4rem;
    }
    .hero p {
        color: #475569;
        font-size: 1.05rem;
        margin-top: 0;
    }

    /* ── Upload card ─────────────────────────── */
    .upload-card {
        background: #ffffff;
        border: 1px solid #e2e8f0;
        border-radius: 16px;
        padding: 2rem;
        margin: 1rem auto;
        max-width: 640px;
        box-shadow: 0 4px 16px rgba(0,0,0,0.07);
    }

    /* ── Slide cards ─────────────────────────── */
    .slide-card {
        background: #ffffff;
        border: 1px solid #e2e8f0;
        border-radius: 12px;
        padding: 1.2rem 1.5rem;
        margin-bottom: 0.8rem;
        box-shadow: 0 2px 8px rgba(0,0,0,0.05);
    }
    .slide-card h3 {
        color: #7c3aed;
        margin: 0 0 0.5rem;
        font-size: 1rem;
    }
    .slide-card ul {
        color: #334155;
        margin: 0;
        padding-left: 1.2rem;
    }
    .slide-card ul li {
        margin-bottom: 0.2rem;
        font-size: 0.92rem;
    }

    /* ── Stats badge row ─────────────────────── */
    .stat-row {
        display: flex;
        gap: 1rem;
        justify-content: center;
        flex-wrap: wrap;
        margin: 1.2rem 0;
    }
    .stat-badge {
        background: #ede9fe;
        border: 1px solid #c4b5fd;
        border-radius: 999px;
        padding: 0.4rem 1.1rem;
        color: #6d28d9;
        font-size: 0.88rem;
        font-weight: 600;
    }

    /* ── Buttons ─────────────────────────────── */
    .stDownloadButton > button {
        background: linear-gradient(135deg, #7c3aed, #2563eb) !important;
        color: white !important;
        border: none !important;
        border-radius: 10px !important;
        padding: 0.6rem 1.6rem !important;
        font-weight: 600 !important;
        font-size: 1rem !important;
        width: 100%;
        transition: opacity 0.2s;
    }
    .stDownloadButton > button:hover {
        opacity: 0.88 !important;
    }

    /* ── File uploader ───────────────────────── */
    [data-testid="stFileUploadDropzone"] {
        background: #f8fafc !important;
        border: 2px dashed #a78bfa !important;
        border-radius: 12px !important;
        color: #475569 !important;
    }

    /* ── Tabs ────────────────────────────────── */
    .stTabs [data-baseweb="tab-list"] {
        background: #e2e8f0;
        border-radius: 10px;
        gap: 4px;
        padding: 4px;
    }
    .stTabs [data-baseweb="tab"] {
        color: #475569 !important;
        border-radius: 8px;
        font-weight: 500;
    }
    .stTabs [aria-selected="true"] {
        background: #ffffff !important;
        color: #7c3aed !important;
        box-shadow: 0 1px 4px rgba(0,0,0,0.10);
    }

    /* ── Code / text area ────────────────────── */
    .stTextArea textarea {
        background: #f8fafc !important;
        color: #1e293b !important;
        border: 1px solid #cbd5e1 !important;
        border-radius: 10px !important;
        font-family: 'Fira Mono', monospace;
        font-size: 0.85rem;
    }

    /* ── Hide default Streamlit chrome ───────── */
    #MainMenu, footer { visibility: hidden; }
    </style>
    """,
    unsafe_allow_html=True,
)

# ─────────────────────────────────────────────
# Core conversion logic
# ─────────────────────────────────────────────

def _table_to_md(rows: list[list[str]]) -> str:
    """Convert a list of row-lists into a Markdown pipe table string."""
    if not rows:
        return ""
    def escape(cell: str) -> str:
        return cell.replace("|", "\\|").replace("\n", " ")
    header = "| " + " | ".join(escape(c) for c in rows[0]) + " |"
    sep    = "| " + " | ".join(["---"] * len(rows[0])) + " |"
    body   = ["| " + " | ".join(escape(c) for c in row) + " |" for row in rows[1:]]
    return "\n".join([header, sep] + body)


def parse_pptx(file_bytes: bytes) -> list[dict]:
    """
    Return a list of slide dicts.
    Each slide: {index, title, content}
    content is an ordered list of items:
      {"type": "text",  "lines": [str, ...]}
      {"type": "table", "rows":  [[str, ...], ...]}
    """
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        content: list[dict] = []

        for shape in slide.shapes:
            # ── Tables ──────────────────────────────
            if shape.has_table:
                rows_data = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                # Skip fully-empty tables
                if any(any(c for c in row) for row in rows_data):
                    content.append({"type": "table", "rows": rows_data})
                continue

            # ── Text shapes ─────────────────────────
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue

            if title is None:
                title = text          # First non-empty text → slide title
            else:
                lines = [sub.strip() for sub in text.split("\n") if sub.strip()]
                if lines:
                    content.append({"type": "text", "lines": lines})

        slides.append({"index": i, "title": title or "Untitled", "content": content})
    return slides


def slides_to_markdown(slides: list[dict], pptx_name: str = "Presentation") -> str:
    """Convert parsed slides to a Markdown string (text + tables)."""
    stem = os.path.splitext(pptx_name)[0]
    lines = [f"# {stem}\n"]
    for s in slides:
        lines.append(f"## Slide {s['index']}: {s['title']}")
        for item in s["content"]:
            if item["type"] == "text":
                for line in item["lines"]:
                    lines.append(f"- {line}")
            elif item["type"] == "table":
                lines.append("")                       # blank line before table
                lines.append(_table_to_md(item["rows"]))
                lines.append("")                       # blank line after table
        lines.append("")
    return "\n".join(lines)


# ─────────────────────────────────────────────
# Excel conversion logic
# ─────────────────────────────────────────────

def _split_into_blocks(df: pd.DataFrame) -> list[list[list[str]]]:
    """
    Split a sheet DataFrame into table blocks separated by fully-blank rows.
    Returns a list of blocks; each block is a list of rows (list of str).
    """
    blocks: list[list[list[str]]] = []
    current: list[list[str]] = []
    for _, row in df.iterrows():
        row_vals = [str(v).strip() if str(v).strip() not in ("", "nan") else "" for v in row]
        if all(v == "" for v in row_vals):
            if current:
                blocks.append(current)
                current = []
        else:
            current.append(row_vals)
    if current:
        blocks.append(current)
    return blocks


def parse_excel(file_bytes: bytes) -> list[dict]:
    """
    Return a list of sheet dicts: {name, blocks}
    Each block is a list of row-lists (first row = header).
    """
    xl = pd.ExcelFile(io.BytesIO(file_bytes), engine="openpyxl")
    sheets = []
    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name, header=None)
        blocks = _split_into_blocks(df)
        # Filter out blocks that are entirely empty strings
        blocks = [b for b in blocks if any(any(c for c in r) for r in b)]
        sheets.append({"name": sheet_name, "blocks": blocks})
    return sheets


def excel_to_markdown(sheets: list[dict], xlsx_name: str = "Spreadsheet") -> str:
    """Convert parsed Excel sheets to Markdown with one table per block."""
    stem = os.path.splitext(xlsx_name)[0]
    lines = [f"# {stem}\n"]
    for sheet in sheets:
        lines.append(f"## Sheet: {sheet['name']}")
        if not sheet["blocks"]:
            lines.append("_No data found on this sheet._")
            lines.append("")
            continue
        for block in sheet["blocks"]:
            # Normalise row widths
            max_cols = max(len(r) for r in block)
            padded = [r + [""] * (max_cols - len(r)) for r in block]
            lines.append("")
            lines.append(_table_to_md(padded))
            lines.append("")
        lines.append("")
    return "\n".join(lines)


# ─────────────────────────────────────────────
# Word conversion logic
# ─────────────────────────────────────────────

_HEADING_PREFIX = {
    "Heading 1": "#",
    "Heading 2": "##",
    "Heading 3": "###",
    "Heading 4": "####",
    "Heading 5": "#####",
    "Heading 6": "######",
}
_BULLET_STYLES  = {"List Bullet", "List Bullet 2", "List Bullet 3", "List Paragraph"}
_NUMBER_STYLES  = {"List Number", "List Number 2", "List Number 3"}


def _render_runs(para: DocxParagraph) -> str:
    """Render paragraph runs with bold / italic Markdown markers."""
    parts = []
    for run in para.runs:
        t = run.text
        if not t:
            continue
        if run.bold and run.italic:
            t = f"***{t}***"
        elif run.bold:
            t = f"**{t}**"
        elif run.italic:
            t = f"*{t}*"
        parts.append(t)
    return "".join(parts).strip()


def parse_docx(file_bytes: bytes) -> list[dict]:
    """
    Walk the Word document body in element order (preserving
    paragraph/table interleaving) and return a list of items:
      {"type": "heading",   "level": int, "text": str}
      {"type": "bullet",   "text": str}
      {"type": "number",   "text": str}
      {"type": "paragraph","text": str}
      {"type": "table",    "rows": [[str, ...], ...]}
    """
    doc   = Document(io.BytesIO(file_bytes))
    items: list[dict] = []

    for element in doc.element.body:
        tag = element.tag

        # ── Table ────────────────────────────────
        if tag == qn("w:tbl"):
            tbl  = DocxTable(element, doc)
            rows = [
                [" ".join(cell.text.split()) for cell in row.cells]
                for row in tbl.rows
            ]
            # de-duplicate merged cells that repeat adjacent text
            deduped = []
            for row in rows:
                clean = []
                prev  = None
                for cell in row:
                    clean.append(cell if cell != prev else "")
                    prev = cell
                deduped.append(clean)
            if any(any(c for c in r) for r in deduped):
                items.append({"type": "table", "rows": deduped})
            continue

        # ── Paragraph ───────────────────────────
        if tag == qn("w:p"):
            para  = DocxParagraph(element, doc)
            style = para.style.name if para.style else "Normal"
            text  = _render_runs(para)
            if not text:
                continue

            if style in _HEADING_PREFIX:
                items.append({"type": "heading",
                               "level": int(style[-1]),
                               "text": text})
            elif style in _BULLET_STYLES:
                items.append({"type": "bullet", "text": text})
            elif style in _NUMBER_STYLES:
                items.append({"type": "number", "text": text})
            else:
                items.append({"type": "paragraph", "text": text})

    return items


def docx_to_markdown(items: list[dict], docx_name: str = "Document") -> str:
    """Convert parsed Word items to a Markdown string."""
    stem  = os.path.splitext(docx_name)[0]
    lines = [f"# {stem}\n"]
    num_counter = 0   # track numbered-list counter

    for item in items:
        t = item["type"]

        if t == "heading":
            num_counter = 0
            prefix = _HEADING_PREFIX.get(f"Heading {item['level']}", "##")
            lines.append(f"{prefix} {item['text']}")

        elif t == "bullet":
            num_counter = 0
            lines.append(f"- {item['text']}")

        elif t == "number":
            num_counter += 1
            lines.append(f"{num_counter}. {item['text']}")

        elif t == "paragraph":
            num_counter = 0
            lines.append(item["text"])
            lines.append("")   # blank line after prose paragraph

        elif t == "table":
            num_counter = 0
            lines.append("")
            lines.append(_table_to_md(item["rows"]))
            lines.append("")

    return "\n".join(lines)


# ─────────────────────────────────────────────
# PDF conversion logic
# ─────────────────────────────────────────────

def parse_pdf(file_bytes: bytes) -> list[dict]:
    """
    Parse each PDF page into an ordered list of items preserving
    reading order by sorting text + table blocks by their y-position.

    Returns a list of page dicts:
      {"index": int, "items": [{"type": "text"|"table", ...}]}

    Returns an empty list for image-only (scanned) PDFs.
    """
    pages = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            content: list[tuple[float, dict]] = []  # (y_pos, item)

            # ── Detect tables and their bounding boxes ────────
            pdf_tables  = page.find_tables()
            table_bboxes = [t.bbox for t in pdf_tables]  # (x0, top, x1, bottom)

            # ── Extract words NOT inside any table bbox ────────
            def _in_table(word: dict) -> bool:
                wx0, wy0 = word["x0"], word["top"]
                wx1, wy1 = word["x1"], word["bottom"]
                for (tx0, ty0, tx1, ty1) in table_bboxes:
                    if wx0 >= tx0 - 2 and wy0 >= ty0 - 2 and wx1 <= tx1 + 2 and wy1 <= ty1 + 2:
                        return True
                return False

            words = page.extract_words() or []
            non_table_words = [w for w in words if not _in_table(w)]

            # Group non-table words into lines by approximate y-position
            line_map: dict[int, list[str]] = {}
            for w in non_table_words:
                y_key = round(w["top"] / 3) * 3  # bucket every ~3 pts
                line_map.setdefault(y_key, []).append(w["text"])

            for y, line_words in line_map.items():
                line_text = " ".join(line_words).strip()
                if line_text:
                    content.append((float(y), {"type": "text", "text": line_text}))

            # ── Add table items with their top-y position ────────
            for pdf_table in pdf_tables:
                table_y = pdf_table.bbox[1]          # top edge
                rows = pdf_table.extract() or []
                clean_rows = [
                    [str(cell or "").strip().replace("\n", " ") for cell in row]
                    for row in rows
                ]
                clean_rows = [r for r in clean_rows if any(c for c in r)]
                if clean_rows:
                    content.append((table_y, {"type": "table", "rows": clean_rows}))

            # Sort by y to restore reading order
            content.sort(key=lambda x: x[0])
            pages.append({"index": i, "items": [item for _, item in content]})

    return pages


def pdf_to_markdown(pages: list[dict], pdf_name: str = "Document") -> str:
    """Convert parsed PDF pages to a Markdown string."""
    stem  = os.path.splitext(pdf_name)[0]
    lines = [f"# {stem}\n"]
    for page in pages:
        lines.append(f"## Page {page['index']}")
        prev_was_text = False
        for item in page["items"]:
            if item["type"] == "text":
                lines.append(item["text"])
                prev_was_text = True
            elif item["type"] == "table":
                lines.append("")                          # blank before table
                lines.append(_table_to_md(item["rows"]))
                lines.append("")                          # blank after table
                prev_was_text = False
        if prev_was_text:
            lines.append("")   # blank line between pages
        lines.append("")
    return "\n".join(lines)

# ─────────────────────────────────────────────
# UI
# ─────────────────────────────────────────────

st.markdown(
    """
    <div class="hero">
        <h1>📑 Docs → Markdown</h1>
        <p>Upload a <strong>.pptx</strong>, <strong>.xlsx</strong>, <strong>.docx</strong>, or <strong>.pdf</strong> file and get a clean Markdown file instantly.</p>
    </div>
    """,
    unsafe_allow_html=True,
)

# ── Upload ────────────────────────────
st.markdown('<div class="upload-card">', unsafe_allow_html=True)
uploaded = st.file_uploader(
    "Drop a .pptx, .xlsx, .docx, or .pdf file here, or click to browse",
    type=["pptx", "xlsx", "xls", "docx", "pdf"],
    label_visibility="collapsed",
)
st.markdown("</div>", unsafe_allow_html=True)

# ── Helpers shared by both renderers ────────
def _html_table(rows: list[list[str]]) -> str:
    """Render a list-of-rows as an HTML table string."""
    if not rows:
        return ""
    thead = "<tr>" + "".join(
        f"<th style='padding:6px 12px;border-bottom:2px solid #7c3aed;color:#7c3aed;text-align:left'>{c}</th>"
        for c in rows[0]
    ) + "</tr>"
    tbody = "".join(
        "<tr>" + "".join(
            f"<td style='padding:5px 12px;border-bottom:1px solid #e2e8f0;color:#334155'>{c}</td>"
            for c in row
        ) + "</tr>"
        for row in rows[1:]
    )
    return (
        f"<table style='width:100%;border-collapse:collapse;margin:0.6rem 0;font-size:0.88rem'>"
        f"<thead>{thead}</thead><tbody>{tbody}</tbody></table>"
    )


# ── Process ────────────────────────────────
if uploaded is not None:
    file_bytes  = uploaded.read()
    file_name   = uploaded.name
    ext         = os.path.splitext(file_name)[1].lower()
    md_filename = os.path.splitext(file_name)[0] + ".md"

    # ── PowerPoint ───────────────────────────
    if ext == ".pptx":
        with st.spinner("Converting slides…"):
            slides  = parse_pptx(file_bytes)
            md_text = slides_to_markdown(slides, file_name)

        slide_count  = len(slides)
        bullet_count = sum(
            len(item["lines"]) for s in slides
            for item in s["content"] if item["type"] == "text"
        )
        table_count = sum(
            1 for s in slides
            for item in s["content"] if item["type"] == "table"
        )
        word_count = len(md_text.split())

        table_badge = (
            f'<span class="stat-badge">📊 {table_count} table{"s" if table_count != 1 else ""}</span>'
            if table_count else ""
        )
        st.markdown(
            f"""
            <div class="stat-row">
                <span class="stat-badge">🗂 {slide_count} slides</span>
                <span class="stat-badge">• {bullet_count} bullet points</span>
                {table_badge}
                <span class="stat-badge">📝 ~{word_count} words</span>
            </div>
            """,
            unsafe_allow_html=True,
        )

        col_dl, col_gap = st.columns([1, 2])
        with col_dl:
            st.download_button(
                label="⬇️  Download Markdown",
                data=md_text.encode("utf-8"),
                file_name=md_filename,
                mime="text/markdown",
            )

        st.divider()

        tab_preview, tab_raw, tab_detail = st.tabs(
            ["🖼 Rendered Preview", "📄 Raw Markdown", "🗂 Slide-by-Slide"]
        )
        with tab_preview:
            st.markdown(md_text)
        with tab_raw:
            st.text_area("Markdown source", value=md_text, height=520, label_visibility="collapsed")
        with tab_detail:
            for s in slides:
                inner_html = ""
                for item in s["content"]:
                    if item["type"] == "text":
                        lis = "".join(f"<li>{ln}</li>" for ln in item["lines"])
                        inner_html += f"<ul style='color:#334155;padding-left:1.2rem;margin:0.4rem 0'>{lis}</ul>"
                    elif item["type"] == "table":
                        inner_html += _html_table(item["rows"])
                if not inner_html:
                    inner_html = "<em style='color:#94a3b8'>No text or table content</em>"
                st.markdown(
                    f"<div class='slide-card'><h3>Slide {s['index']}: {s['title']}</h3>{inner_html}</div>",
                    unsafe_allow_html=True,
                )

    # ── Excel ────────────────────────────────
    elif ext in (".xlsx", ".xls"):
        with st.spinner("Reading sheets…"):
            sheets  = parse_excel(file_bytes)
            md_text = excel_to_markdown(sheets, file_name)

        sheet_count = len(sheets)
        table_count = sum(len(s["blocks"]) for s in sheets)
        word_count  = len(md_text.split())

        st.markdown(
            f"""
            <div class="stat-row">
                <span class="stat-badge">📄 {sheet_count} sheet{"s" if sheet_count != 1 else ""}</span>
                <span class="stat-badge">📊 {table_count} table block{"s" if table_count != 1 else ""}</span>
                <span class="stat-badge">📝 ~{word_count} words</span>
            </div>
            """,
            unsafe_allow_html=True,
        )

        col_dl, col_gap = st.columns([1, 2])
        with col_dl:
            st.download_button(
                label="⬇️  Download Markdown",
                data=md_text.encode("utf-8"),
                file_name=md_filename,
                mime="text/markdown",
            )

        st.divider()

        tab_preview, tab_raw, tab_detail = st.tabs(
            ["🖼 Rendered Preview", "📄 Raw Markdown", "📄 Sheet-by-Sheet"]
        )
        with tab_preview:
            st.markdown(md_text)
        with tab_raw:
            st.text_area("Markdown source", value=md_text, height=520, label_visibility="collapsed")
        with tab_detail:
            for sheet in sheets:
                header_html = f"<h3 style='color:#7c3aed;margin-bottom:0.5rem'>📄 {sheet['name']}</h3>"
                if not sheet["blocks"]:
                    body_html = "<em style='color:#94a3b8'>No data found on this sheet.</em>"
                else:
                    body_html = "".join(_html_table(b) for b in sheet["blocks"])
                st.markdown(
                    f"<div class='slide-card'>{header_html}{body_html}</div>",
                    unsafe_allow_html=True,
                )

    # ── Word ─────────────────────────────────────────────────────
    elif ext == ".docx":
        with st.spinner("Parsing document…"):
            doc_items = parse_docx(file_bytes)
            md_text   = docx_to_markdown(doc_items, file_name)

        heading_count = sum(1 for i in doc_items if i["type"] == "heading")
        table_count   = sum(1 for i in doc_items if i["type"] == "table")
        bullet_count  = sum(1 for i in doc_items if i["type"] in ("bullet", "number"))
        word_count    = len(md_text.split())
        table_badge   = (
            f'<span class="stat-badge">📊 {table_count} table{"s" if table_count != 1 else ""}</span>'
            if table_count else ""
        )
        st.markdown(
            f"""
            <div class="stat-row">
                <span class="stat-badge">🏷️ {heading_count} heading{"s" if heading_count != 1 else ""}</span>
                <span class="stat-badge">• {bullet_count} list items</span>
                {table_badge}
                <span class="stat-badge">📝 ~{word_count} words</span>
            </div>
            """,
            unsafe_allow_html=True,
        )
        col_dl, col_gap = st.columns([1, 2])
        with col_dl:
            st.download_button("⬇️  Download Markdown", md_text.encode(), md_filename, "text/markdown")
        st.divider()
        tab_preview, tab_raw, tab_detail = st.tabs(
            ["🖼 Rendered Preview", "📄 Raw Markdown", "📋 Document Structure"]
        )
        with tab_preview:
            st.markdown(md_text)
        with tab_raw:
            st.text_area("Markdown source", value=md_text, height=520, label_visibility="collapsed")
        with tab_detail:
            for item in doc_items:
                t = item["type"]
                if t == "heading":
                    prefix = "#" * item["level"]
                    st.markdown(
                        f"<div class='slide-card'><h3 style='color:#7c3aed;margin:0'>{prefix} {item['text']}</h3></div>",
                        unsafe_allow_html=True,
                    )
                elif t == "bullet":
                    st.markdown(f"&nbsp;&nbsp;&nbsp;&nbsp;• {item['text']}")
                elif t == "number":
                    st.markdown(f"&nbsp;&nbsp;&nbsp;&nbsp;# {item['text']}")
                elif t == "table":
                    st.markdown(
                        f"<div class='slide-card'>{_html_table(item['rows'])}</div>",
                        unsafe_allow_html=True,
                    )
                else:
                    st.markdown(item["text"])

    # ── PDF ─────────────────────────────────
    elif ext == ".pdf":
        with st.spinner("Extracting PDF content…"):
            pdf_pages = parse_pdf(file_bytes)
            md_text   = pdf_to_markdown(pdf_pages, file_name)

        page_count  = len(pdf_pages)
        table_count = sum(1 for p in pdf_pages for i in p["items"] if i["type"] == "table")
        word_count  = len(md_text.split())
        has_text    = any(i["type"] == "text" for p in pdf_pages for i in p["items"])

        if not has_text:
            st.warning(
                "⚠️ This PDF appears to be image-based (scanned). "
                "Text extraction requires OCR, which is not yet supported. "
                "The downloaded Markdown will contain only page headings.",
                icon="📷",
            )

        table_badge = (
            f'<span class="stat-badge">📊 {table_count} table{"s" if table_count != 1 else ""}</span>'
            if table_count else ""
        )
        st.markdown(
            f"""
            <div class="stat-row">
                <span class="stat-badge">📄 {page_count} page{"s" if page_count != 1 else ""}</span>
                {table_badge}
                <span class="stat-badge">📝 ~{word_count} words</span>
            </div>
            """,
            unsafe_allow_html=True,
        )
        col_dl, col_gap = st.columns([1, 2])
        with col_dl:
            st.download_button("⬇️  Download Markdown", md_text.encode(), md_filename, "text/markdown")
        st.divider()
        tab_preview, tab_raw, tab_detail = st.tabs(
            ["🖼 Rendered Preview", "📄 Raw Markdown", "📄 Page-by-Page"]
        )
        with tab_preview:
            st.markdown(md_text)
        with tab_raw:
            st.text_area("Markdown source", value=md_text, height=520, label_visibility="collapsed")
        with tab_detail:
            for page in pdf_pages:
                inner_html = ""
                for item in page["items"]:
                    if item["type"] == "text":
                        inner_html += f"<p style='color:#334155;margin:0.2rem 0;font-size:0.9rem'>{item['text']}</p>"
                    elif item["type"] == "table":
                        inner_html += _html_table(item["rows"])
                if not inner_html:
                    inner_html = "<em style='color:#94a3b8'>No extractable text on this page (may be an image).</em>"
                st.markdown(
                    f"<div class='slide-card'><h3 style='color:#7c3aed;margin:0 0 0.5rem'>Page {page['index']}</h3>{inner_html}</div>",
                    unsafe_allow_html=True,
                )

else:
    st.info(
        "👆 Upload a `.pptx`, `.xlsx`, `.docx`, or `.pdf` file above to get started.",
        icon="💡",
    )
